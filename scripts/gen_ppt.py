#!/usr/bin/env python3
"""
IPD 项目评审 PPT 生成脚本模板

使用方法：
    python gen_ppt.py --stage charter --project "项目名" --type basic
    python gen_ppt.py --stage pdcp --project "项目名" --type applied
    python gen_ppt.py --stage adcp --project "项目名" --type basic
    python gen_ppt.py --stage transfer --project "项目名" --type applied

参数：
    --stage     阶段：charter / pdcp / adcp / transfer
    --project   项目名称
    --type      项目类型：basic（基础研究）/ applied（应用研究）
    --product   产品线：meiling（美菱）/ ac（空调），默认 meiling
    --output    输出路径，默认当前目录
"""

import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


# PPT 模板路径
TEMPLATE_PATH = os.path.expanduser(
    "~/SynologyDrive/虹美公司/部门/美菱PPT模板.pptx"
)

# 人力成本基准
COST_PER_PERSON_DAY = 970


def create_presentation(project_name, stage, project_type, product_line="meiling"):
    """创建评审 PPT"""
    prs = Presentation(TEMPLATE_PATH)

    # 添加封面
    add_cover_slide(prs, project_name, stage)

    # 添加目录
    add_toc_slide(prs, stage)

    # 根据阶段添加内容页
    if stage == "charter":
        add_charter_slides(prs, project_name, project_type)
    elif stage == "pdcp":
        add_pdcp_slides(prs, project_name, project_type)
    elif stage == "adcp":
        add_adcp_slides(prs, project_name, project_type)
    elif stage == "transfer":
        add_transfer_slides(prs, project_name, project_type)

    # 添加结尾页
    add_ending_slide(prs)

    return prs


def add_cover_slide(prs, project_name, stage):
    """添加封面"""
    stage_names = {
        "charter": "立项（Charter）汇报材料",
        "pdcp": "总体设计方案评审（PDCP）汇报材料",
        "adcp": "可获得性评审（ADCP）汇报材料",
        "transfer": "转移阶段汇报材料",
    }
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = project_name
    slide.placeholders[1].text = stage_names.get(stage, "汇报材料")


def add_toc_slide(prs, stage):
    """添加目录页"""
    toc_items = {
        "charter": [
            "一、项目背景",
            "二、研究内容",
            "三、项目价值",
            "四、项目目标（Q/C/D）",
            "五、关键资源计划",
            "六、风险及问题管理",
            "七、财务计划",
            "八、结论",
        ],
        "pdcp": [
            "一、项目简介",
            "二、项目技术目标",
            "三、项目研究方法",
            "四、技术方案及关键点",
            "五、关键资源计划",
            "六、项目风险管理",
            "七、知识产权方案",
            "八、环境及职业健康影响",
            "九、项目成员构成",
        ],
        "adcp": [
            "一、项目简介",
            "二、项目目标完成情况",
            "三、项目转移计划",
            "四、关键资源计划",
            "五、项目转移风险管理",
            "六、项目技术创新总结",
            "七、经济效益和应用前景",
            "八、项目过程运作",
        ],
        "transfer": [
            "一、项目简介",
            "二、项目目标转移情况",
            "三、项目技术成果",
            "四、技术创新点推广应用",
            "五、转移过程运作",
        ],
    }
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.placeholders[0].text = "目录"
    body = slide.placeholders[1]
    for i, item in enumerate(toc_items.get(stage, [])):
        if i == 0:
            body.text_frame.paragraphs[0].text = item
        else:
            p = body.text_frame.add_paragraph()
            p.text = item


def add_content_slide(prs, title, points, notes=""):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.placeholders[0].text = title
    body = slide.placeholders[1]
    for i, point in enumerate(points):
        if i == 0:
            body.text_frame.paragraphs[0].text = point
        else:
            p = body.text_frame.add_paragraph()
            p.text = point
            p.level = 0
    # 添加备注
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


def add_charter_slides(prs, project_name, project_type):
    """添加 Charter 阶段内容页"""
    # 项目背景
    add_content_slide(
        prs,
        "一、项目背景",
        [
            "1. 行业趋势：[待填写]",
            "2. 技术现状：[待填写]",
            "3. 痛点问题：[待填写]",
            "4. 需求来源：[待填写]",
        ],
        "讲解要点：说明为什么是现在做这件事，用数据支撑行业趋势",
    )

    # 研究内容
    add_content_slide(
        prs,
        "二、研究内容",
        [
            "1. 核心研究方向：[待填写]",
            "2. 技术路线概述：[待填写]",
            "3. 关键方法论：[待填写]",
        ],
        "讲解要点：概述技术路线，不需要展开到实现细节",
    )

    # 项目价值（重点页）
    add_content_slide(
        prs,
        "三、项目价值",
        [
            "技术价值：[待填写]",
            "  - 突破了什么技术瓶颈",
            "  - 建立了什么能力",
            "经济价值：[待填写]",
            "  - 预期经济效益",
            "  - 投入产出比",
            "应用价值：[待填写]",
            "  - 可应用的产品/场景",
            "  - 市场前景",
        ],
        "讲解要点：这是最重要的部分！用数据和案例支撑价值论证",
    )

    # 项目目标
    add_content_slide(
        prs,
        "四、项目目标（Q/C/D）",
        [
            "Q 指标（技术规格）：[待填写]",
            "C 指标（财务）：总预算 [待填写] 万元",
            "D 指标（进度）：[待填写] 立项 → [待填写] 验收",
        ],
        "讲解要点：明确、可量化的目标",
    )

    # 关键资源
    add_content_slide(
        prs,
        "五、关键资源计划",
        [
            "人力：[待填写] 人天 × 970 = [自动计算] 万元",
            "设备：[待填写]",
            "外部资源：[待填写]",
        ],
    )

    # 风险管理
    add_content_slide(
        prs,
        "六、风险及问题管理",
        [
            "风险1：[待填写]（高/中/低）",
            "  应对措施：[待填写]",
            "风险2：[待填写]（高/中/低）",
            "  应对措施：[待填写]",
        ],
    )

    # 财务计划
    add_content_slide(
        prs,
        "七、财务计划",
        [
            "人力成本：[待填写] 万元",
            "服务器：[待填写] 万元",
            "试制费：[待填写] 万元",
            "其他：[待填写] 万元",
            "合计：[自动计算] 万元",
        ],
    )

    # 结论
    add_content_slide(
        prs,
        "八、结论",
        [
            "1. 项目必要性：[一句话总结]",
            "2. 预期成果：[一句话总结]",
            "3. 下一步行动：[一句话总结]",
        ],
        "讲解要点：用 2-3 句话有力地总结，给评审委员留下深刻印象",
    )


def add_pdcp_slides(prs, project_name, project_type):
    """添加 PDCP 阶段内容页"""
    slides_content = [
        ("一、项目简介", ["项目背景：[待填写]", "项目目标：[待填写]", "项目范围：[待填写]"]),
        ("二、项目技术目标", ["Q 指标细化：[待填写]", "C 指标各阶段分解：[待填写]", "D 指标里程碑：[待填写]"]),
        ("三、项目研究方法", ["技术路线图：[待填写]", "分步实施计划：[待填写]", "各阶段产出：[待填写]"]),
        ("四、技术方案及关键点", ["系统架构：[待填写]", "核心算法/模型：[待填写]", "关键技术难点：[待填写]", "可行性分析：[待填写]"]),
        ("五、关键资源计划", ["人力：[待填写]", "设备：[待填写]", "关键技术可获得性：[待填写]"]),
        ("六、项目风险管理", ["风险1：[待填写]", "风险2：[待填写]"]),
        ("七、知识产权方案", ["专利检索：[待填写]", "专利规划：[待填写]"]),
        ("八、环境及职业健康影响", ["影响分析：[待填写]"]),
        ("九、项目成员构成", ["项目经理：[待填写]", "核心成员：[待填写]"]),
    ]
    for title, points in slides_content:
        add_content_slide(prs, title, points)


def add_adcp_slides(prs, project_name, project_type):
    """添加 ADCP 阶段内容页"""
    slides_content = [
        ("一、项目简介", ["项目概述：[待填写]"]),
        ("二、项目目标完成情况", [
            "Q 指标偏差分析：",
            "  指标1：目标 [待填写] → 实际 [待填写]",
            "C 指标偏差分析：",
            "  计划 [待填写] 万 → 实际 [待填写] 万",
            "D 指标偏差分析：",
            "  计划 [待填写] → 实际 [待填写]",
        ]),
        ("三、项目转移计划", ["应用产品/平台：[待填写]", "转移条件：[待填写]", "完成时间：[待填写]"]),
        ("四、关键资源计划", ["转移阶段资源：[待填写]"]),
        ("五、项目转移风险管理", ["风险：[待填写]", "应对：[待填写]"]),
        ("六、项目技术创新总结", ["创新点1：[待填写]", "创新点2：[待填写]", "专利/论文：[待填写]"]),
        ("七、经济效益和应用前景", ["已实现效益：[待填写]", "市场规模：[待填写]", "推广策略：[待填写]"]),
        ("八、项目过程运作", ["亮点：[待填写]", "不足：[待填写]", "经验教训：[待填写]"]),
    ]
    for title, points in slides_content:
        add_content_slide(prs, title, points)


def add_transfer_slides(prs, project_name, project_type):
    """添加转移阶段内容页"""
    slides_content = [
        ("一、项目简介", ["项目基本信息：[待填写]"]),
        ("二、项目目标转移情况", [
            "Q 指标转移：",
            "  正式样机测试数据：[待填写]",
            "  正式样机评审：[待填写]",
            "  小批试制：[待填写]",
            "  试生产评审：[待填写]",
            "C 指标转移：[待填写]",
            "D 指标转移：[待填写]",
        ]),
        ("三、项目技术成果", ["方法/规范/标准：[待填写]", "专利：[待填写]", "论文：[待填写]"]),
        ("四、技术创新点推广应用", ["推广产品/平台：[待填写]", "推广效果：[待填写]", "后续计划：[待填写]"]),
        ("五、转移过程运作", ["亮点：[待填写]", "不足：[待填写]", "经验教训：[待填写]"]),
    ]
    for title, points in slides_content:
        add_content_slide(prs, title, points)


def add_ending_slide(prs):
    """添加结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = "谢谢！"
    slide.placeholders[1].text = "请各位领导和专家指导"


def main():
    parser = argparse.ArgumentParser(description="IPD 项目评审 PPT 生成")
    parser.add_argument("--stage", required=True, choices=["charter", "pdcp", "adcp", "transfer"])
    parser.add_argument("--project", required=True, help="项目名称")
    parser.add_argument("--type", default="basic", choices=["basic", "applied"])
    parser.add_argument("--product", default="meiling", choices=["meiling", "ac"])
    parser.add_argument("--output", default=".", help="输出路径")
    args = parser.parse_args()

    prs = create_presentation(args.project, args.stage, args.type, args.product)

    filename = f"{args.project}-{args.stage.upper()}-汇报材料.pptx"
    output_path = os.path.join(args.output, filename)
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")


if __name__ == "__main__":
    main()
