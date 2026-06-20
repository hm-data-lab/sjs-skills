#!/usr/bin/env python3
"""
IPD 项目评审 PPT 生成脚本

基于内置模板，按页面规则填充内容。
模板母版不动，只改内容：封面替换信息、内容页替换子标题、表格页清空填入数据。

使用方法：
    python gen_ppt.py --outline 项目名-CHARTER-v2.md --product-line ml --output ./PPT/

参数：
    --outline       确认版大纲 Markdown 文件路径（必需）
    --product-line  产品线：hm（虹美/空调/通用平台）或 ml（美菱/冰箱/洗衣机/小电日电/数字化）
    --output        输出目录，默认当前目录
    --stage         阶段（可选，默认从文件名解析）
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu


def load_config() -> dict:
    """尝试加载 config.json"""
    script_dir = Path(__file__).resolve().parent
    skill_dir = script_dir.parent
    config_path = skill_dir / "config.json"
    if config_path.exists():
        with open(config_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def resolve_template(product_line: str, stage: str) -> str:
    """根据阶段解析内置模板路径（通用模板）"""
    script_dir = Path(__file__).resolve().parent
    assets_dir = script_dir.parent / "assets"
    template_name = f"{stage}-template.pptx"
    template_path = assets_dir / template_name
    if not template_path.exists():
        print(f"错误：内置模板不存在: {template_path}", file=sys.stderr)
        available = [f.name for f in assets_dir.glob("*-template.pptx")]
        print(f"可用模板：{', '.join(available)}", file=sys.stderr)
        sys.exit(1)
    return str(template_path)


def get_company_info(product_line: str) -> dict:
    """根据产品线返回公司名称和承担单位"""
    info = {
        "hm": {"company": "四川虹美智能科技有限公司", "department": "数据技术所"},
        "ml": {"company": "长虹美菱股份有限公司", "department": "AI数智研发部"},
    }
    return info.get(product_line, info["ml"])


# ── 大纲解析 ──────────────────────────────────────────────

def parse_outline(md_path: str) -> dict:
    """解析大纲 Markdown 文件"""
    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    lines = content.split("\n")
    result = {"title": "", "meta": "", "sections": [], "notes_map": {}}

    in_notes_section = False
    current_section = None
    current_slide = None
    last_chart_hint = None

    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if not line:
            i += 1
            continue

        stripped = line.strip()
        if stripped.startswith("<!--"):
            chart_match = re.search(r"chart:\s*(\w+)", stripped)
            if chart_match:
                last_chart_hint = chart_match.group(1)
            i += 1
            continue

        if in_notes_section:
            if stripped.startswith("## ") and not stripped.startswith("### "):
                heading = stripped[3:].strip()
                notes_lines = []
                i += 1
                while i < len(lines):
                    nl = lines[i].rstrip()
                    if nl.startswith("## ") or nl.startswith("# "):
                        break
                    if nl.startswith("- "):
                        notes_lines.append(nl[2:].strip())
                    elif nl and not nl.strip().startswith("<!--"):
                        notes_lines.append(nl.strip())
                    i += 1
                result["notes_map"][heading] = "\n".join(notes_lines)
            else:
                i += 1
            continue

        if stripped.startswith("# ") and not stripped.startswith("## "):
            result["title"] = stripped[2:].strip()
            i += 1
            continue

        if stripped.startswith("> ") and "产品线" in stripped:
            result["meta"] = stripped[2:].strip()
            i += 1
            continue

        if stripped.startswith("## 备注区") or stripped == "## 备注":
            in_notes_section = True
            i += 1
            continue

        if stripped.startswith("## ") and not stripped.startswith("### "):
            heading = stripped[3:].strip()
            current_section = {"heading": heading, "slides": []}
            result["sections"].append(current_section)
            current_slide = None
            last_chart_hint = None
            i += 1
            continue

        if stripped.startswith("### "):
            subtitle = stripped[4:].strip()
            current_slide = {"subtitle": subtitle, "points": [], "chart": last_chart_hint}
            if current_section is not None:
                current_section["slides"].append(current_slide)
            last_chart_hint = None
            i += 1
            continue

        if stripped.startswith("|"):
            if current_slide is None and current_section is not None:
                current_slide = {"subtitle": None, "points": [], "chart": last_chart_hint}
                current_section["slides"].append(current_slide)
                last_chart_hint = None
            table_rows = []
            while i < len(lines):
                tl = lines[i].rstrip().strip()
                if tl.startswith("|"):
                    if not re.match(r"^\|[\s\-:|]+\|$", tl):
                        inner = tl.strip("|")
                        cells = [c.strip() for c in inner.split("|")]
                        table_rows.append(cells)
                    i += 1
                else:
                    break
            if table_rows and current_slide is not None:
                for row in table_rows:
                    current_slide["points"].append(" | ".join(row))
            continue

        if stripped.startswith("- "):
            point = stripped[2:].strip()
            if current_slide is None and current_section is not None:
                current_slide = {"subtitle": None, "points": [], "chart": last_chart_hint}
                current_section["slides"].append(current_slide)
                last_chart_hint = None
            if current_slide is not None:
                current_slide["points"].append(point)
            i += 1
            continue

        if current_slide is not None and not stripped.startswith("#"):
            if stripped and stripped != "...":
                current_slide["points"].append(stripped)

        i += 1

    return result


# ── 模板页面操作 ──────────────────────────────────────────

def get_slide_title(slide) -> str:
    """获取 slide 第一个有文字的 shape"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def normalize(text: str) -> str:
    return re.sub(r"\s+", "", text).strip("  .")


def find_table(slide):
    """找到 slide 中的第一个表格"""
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def find_body_shape(slide):
    """找到 slide 中的正文文本框（跳过表格和标题）"""
    found_title = False
    for shape in slide.shapes:
        if shape.has_table:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and not found_title:
                found_title = True
                continue
            if shape.has_text_frame:
                return shape
    return None


def extract_template_style(tf):
    """从 text frame 中提取模板原始字体样式"""
    style = {"name": None, "size": None, "bold": None, "color": None}
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text.strip():
                style["name"] = run.font.name
                style["size"] = run.font.size
                style["bold"] = run.font.bold
                try:
                    if run.font.color and run.font.color.type is not None:
                        style["color"] = run.font.color.rgb
                except (AttributeError, TypeError):
                    pass
                return style
    return style


def apply_style(run, style):
    """将模板样式应用到 run"""
    if style["name"]:
        run.font.name = style["name"]
    if style["size"]:
        run.font.size = style["size"]
    if style["bold"] is not None:
        run.font.bold = style["bold"]
    if style["color"]:
        run.font.color.rgb = style["color"]


def fix_shape_bounds(shape, prs):
    """修正 shape 的位置和尺寸，确保不超出页面边界"""
    page_width = prs.slide_width
    page_height = prs.slide_height
    max_width = page_width - Emu(60000)
    if shape.width > max_width:
        shape.width = max_width
    if shape.left < 0:
        shape.left = Emu(60000)
    if shape.left + shape.width > page_width:
        shape.width = page_width - shape.left - Emu(30000)
    if shape.top < 0:
        shape.top = Emu(30000)
    if shape.top + shape.height > page_height:
        shape.height = page_height - shape.top - Emu(30000)


def fill_table_data(table, points: list):
    """清空表格数据行（保留表头），填入大纲数据"""
    header_cells = []
    if len(table.rows) > 0:
        header_cells = [table.cell(0, c).text.strip() for c in range(len(table.columns))]

    data_rows = []
    for p in points:
        if " | " in p:
            cells = [c.strip() for c in p.split(" | ")]
            if cells and header_cells:
                c0 = cells[0]
                h0 = header_cells[0]
                if c0 == h0:
                    continue
                if len(c0) < len(h0) and h0.startswith(c0):
                    continue
            data_rows.append(cells)

    if not data_rows:
        return

    # 记录模板字体样式
    style = {"name": None, "size": None, "bold": None}
    if len(table.rows) > 1:
        for para in table.cell(1, 0).text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    style["name"] = run.font.name
                    style["size"] = run.font.size
                    style["bold"] = run.font.bold
                    break
            if style["name"]:
                break

    # 清空数据行
    for row_idx in range(1, len(table.rows)):
        for col_idx in range(len(table.columns)):
            table.cell(row_idx, col_idx).text = ""

    # 填入数据
    template_cols = len(table.columns)
    for target_row, row_data in enumerate(data_rows):
        row_idx = target_row + 1
        if row_idx >= len(table.rows):
            break
        for col_idx in range(min(template_cols, len(row_data))):
            cell = table.cell(row_idx, col_idx)
            value = row_data[col_idx]
            value = re.sub(r"\*\*(.+?)\*\*", r"\1", value)
            cell.text = value
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if style["name"]:
                        run.font.name = style["name"]
                    if style["size"]:
                        run.font.size = style["size"]
                    if style["bold"] is not None:
                        run.font.bold = style["bold"]


# ── 页面类型处理 ──────────────────────────────────────────

def process_cover(slide, outline, prs, product_line="ml"):
    """封面页：替换项目名、类型、日期、公司名、承担单位"""
    title = outline.get("title", "")
    project_name = title.split("—")[0].strip() if "—" in title else title.split("-")[0].strip()
    meta = outline.get("meta", "")
    project_type = "基础研究" if "基础研究" in meta else "应用研究" if "应用研究" in meta else "基础研究"
    company_info = get_company_info(product_line)

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full_text = para.text
                # 替换公司名
                if "长虹美菱股份有限公司" in full_text:
                    full_text = full_text.replace("长虹美菱股份有限公司", company_info["company"])
                if "四川虹美智能科技有限公司" in full_text:
                    full_text = full_text.replace("四川虹美智能科技有限公司", company_info["company"])
                # 替换承担单位
                if "AI数智研发部" in full_text:
                    full_text = full_text.replace("AI数智研发部", company_info["department"])
                if "数据技术所" in full_text:
                    full_text = full_text.replace("数据技术所", company_info["department"])
                # 替换日期
                if "2026年X月" in full_text:
                    full_text = full_text.replace("2026年X月", "2026年6月")
                # 替换项目名
                if "XXXXX" in full_text:
                    full_text = full_text.replace("XXXXX技术研究", project_name)
                    full_text = full_text.replace("XXXXX", project_name)
                # 替换项目类别
                if "基础研究" in full_text and project_type != "基础研究":
                    full_text = full_text.replace("基础研究", project_type)
                if full_text != para.text:
                    for ri, run in enumerate(para.runs):
                        run.text = full_text if ri == 0 else ""


def process_content_slide(slide, md_slides, prs):
    """内容页：Shape [1] 只放子标题，正文内容放在下方空白区域"""
    if not md_slides:
        return

    shapes = list(slide.shapes)
    if len(shapes) < 2:
        return

    subtitle_shape = shapes[1]
    if not subtitle_shape.has_text_frame:
        return

    # 提取模板样式
    style = extract_template_style(subtitle_shape.text_frame)

    md_slide = md_slides[0]
    subtitle = md_slide.get("subtitle", "")
    points = md_slide.get("points", [])

    # 1. 替换 Shape [1] 的子标题文字（只放子标题）
    tf = subtitle_shape.text_frame
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._element.getparent().remove(p._element)

    if subtitle:
        tf.paragraphs[0].text = subtitle
        for run in tf.paragraphs[0].runs:
            apply_style(run, style)
            run.font.bold = True

    # 2. 在子标题下方创建新的 text box 放正文内容
    if points:
        # 计算正文区域位置：子标题下方，留 20pt 间距
        content_top = subtitle_shape.top + subtitle_shape.height + Emu(20000)
        content_left = subtitle_shape.left
        content_width = prs.slide_width - content_left - Emu(60000)  # 右侧留边距
        content_height = prs.slide_height - content_top - Emu(50000)  # 底部留边距

        txBox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        body_tf = txBox.text_frame
        body_tf.word_wrap = True

        first = True
        for point in points:
            point = re.sub(r"\*\*(.+?)\*\*", r"\1", point)
            if first:
                body_tf.paragraphs[0].text = point
                for run in body_tf.paragraphs[0].runs:
                    apply_style(run, style)
                first = False
            else:
                p = body_tf.add_paragraph()
                p.text = point
                for run in p.runs:
                    apply_style(run, style)


def process_table_slide(slide, md_slides, is_scoring=False):
    """表格页：清空数据行，填入大纲数据"""
    if not md_slides:
        return
    table = find_table(slide)
    if table:
        if is_scoring:
            fill_scoring_table(table, md_slides[0]["points"])
        else:
            fill_table_data(table, md_slides[0]["points"])


def fill_scoring_table(table, points):
    """结论评分表：只更新单项得分、总得分、建议等级，保留表头和评分标准"""
    # 解析大纲数据
    data_rows = []
    for p in points:
        if " | " in p:
            cells = [c.strip() for c in p.split(" | ")]
            data_rows.append(cells)

    # 从大纲中提取需要更新的行
    score_row = None  # 单项得分
    total_row = None  # 总得分
    grade_row = None  # 建议项目等级

    for row in data_rows:
        if row[0] == "单项得分":
            score_row = row
        elif row[0] == "总得分":
            total_row = row
        elif row[0] == "建议项目等级":
            grade_row = row

    # 模板行映射：Row 2=单项得分, Row 3=总得分, Row 4=建议等级
    if score_row and len(table.rows) > 2:
        for col_idx in range(1, min(len(table.columns), len(score_row))):
            table.cell(2, col_idx).text = score_row[col_idx]

    if total_row and len(table.rows) > 3:
        # 自动计算总得分
        if score_row and len(score_row) >= 4:
            try:
                importance = float(score_row[1])  # 重要性
                innovation = float(score_row[2])  # 创新先进性
                difficulty = float(score_row[3])  # 难度
                total = importance * 0.4 + innovation * 0.4 + difficulty * 0.2
                total_text = f"重要性*40%+创新先进性*40%+难度*20%={total:.1f}"
                table.cell(3, 1).text = total_text
            except (ValueError, IndexError):
                # 解析失败，保留原文
                for col_idx in range(1, min(len(table.columns), len(total_row))):
                    table.cell(3, col_idx).text = total_row[col_idx]
        else:
            for col_idx in range(1, min(len(table.columns), len(total_row))):
                table.cell(3, col_idx).text = total_row[col_idx]

    if grade_row and len(table.rows) > 4:
        for col_idx in range(1, min(len(table.columns), len(grade_row))):
            table.cell(4, col_idx).text = grade_row[col_idx]


def process_summary_slide(slide, outline):
    """结论-总结页：替换正文和项目经理，保留固定写法"""
    title = outline.get("title", "")
    project_name = title.split("—")[0].strip() if "—" in title else ""

    # 从大纲中提取总结内容（只取 subtitle="总结" 的 slide）
    summary_points = []
    for section in outline["sections"]:
        if "结论" in section["heading"]:
            for s in section["slides"]:
                if s.get("subtitle") == "总结":
                    summary_points = s["points"]
                    break
            if summary_points:
                break

    summary_text = "；".join(summary_points[:3]) if summary_points else ""

    shapes = list(slide.shapes)
    # Shape [2] = 正文描述，Shape [3] = 建议项目经理，Shape [4] = 恳请IPMT

    # 替换 Shape [2]（正文描述）
    if len(shapes) > 2 and shapes[2].has_text_frame and summary_text:
        body = shapes[2]
        style = extract_template_style(body.text_frame)
        tf = body.text_frame
        while len(tf.paragraphs) > 1:
            p = tf.paragraphs[-1]
            p._element.getparent().remove(p._element)
        tf.paragraphs[0].text = summary_text
        for run in tf.paragraphs[0].runs:
            apply_style(run, style)

    # 替换 Shape [3]（建议项目经理）
    if len(shapes) > 3 and shapes[3].has_text_frame:
        for para in shapes[3].text_frame.paragraphs:
            if "建议项目经理为：" in para.text:
                new_text = f"建议项目经理为：{project_name}项目负责人"
                for ri, run in enumerate(para.runs):
                    run.text = new_text if ri == 0 else ""


def process_notes(slide, outline, heading):
    """写入备注"""
    notes = outline["notes_map"].get(heading, "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ── 主流程 ────────────────────────────────────────────────

# 固定不动的 slide 索引
FIXED_SLIDES = {1, 19, 20}  # 目录、专业委员会、结束页

# 封面页
COVER_SLIDE = 0

# 结论-总结页（需要特殊处理）
SUMMARY_SLIDE = 18


def generate_ppt(outline: dict, template_path: str, output_path: str, product_line: str = "ml"):
    """根据大纲和模板生成 PPT"""
    prs = Presentation(template_path)
    slides = list(prs.slides)

    matched = 0

    # 按章节分组大纲内容
    # 每个章节的 sub-slides 按顺序分配给模板中匹配的 slide
    section_slide_counter = {}  # 记录每个章节已处理的 sub-slide 数量

    for slide_num in range(len(slides)):
        if slide_num in FIXED_SLIDES or slide_num == COVER_SLIDE:
            continue

        if slide_num == SUMMARY_SLIDE:
            process_summary_slide(slides[slide_num], outline)
            matched += 1
            continue

        # 找到匹配当前模板 slide 的大纲章节
        slide_title = normalize(get_slide_title(slides[slide_num]))

        matched_section = None
        for section in outline["sections"]:
            heading_norm = normalize(section["heading"])
            if heading_norm in slide_title or slide_title in heading_norm:
                matched_section = section
                break

        if not matched_section:
            continue

        heading = matched_section["heading"]
        md_slides = matched_section["slides"]

        # 获取当前章节已处理的 sub-slide 数量
        if heading not in section_slide_counter:
            section_slide_counter[heading] = 0
        sub_idx = section_slide_counter[heading]

        # 获取对应的 sub-slide
        if sub_idx < len(md_slides):
            current_md = md_slides[sub_idx]
            md_slides_for_page = [current_md]
        else:
            # 没有更多 sub-slide 了，跳过
            continue

        # 判断页面类型并处理
        table = find_table(slides[slide_num])

        if table:
            is_scoring = (slide_num == 17)
            process_table_slide(slides[slide_num], md_slides_for_page, is_scoring=is_scoring)
        elif slide_num == 9:
            # 项目价值页（图表页）- 不处理
            pass
        elif slide_num == 15:
            # 产出收益页 - 替换正文
            body = find_body_shape(slides[slide_num])
            if body and body.has_text_frame:
                style = extract_template_style(body.text_frame)
                tf = body.text_frame
                while len(tf.paragraphs) > 1:
                    p = tf.paragraphs[-1]
                    p._element.getparent().remove(p._element)
                first = True
                for point in current_md["points"]:
                    point = re.sub(r"\*\*(.+?)\*\*", r"\1", point)
                    if first:
                        tf.paragraphs[0].text = point
                        for run in tf.paragraphs[0].runs:
                            apply_style(run, style)
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = point
                        for run in p.runs:
                            apply_style(run, style)
        else:
            process_content_slide(slides[slide_num], md_slides_for_page, prs)

        # 写入备注（只在第一个 sub-slide 时写入）
        if sub_idx == 0:
            process_notes(slides[slide_num], outline, heading)

        section_slide_counter[heading] += 1
        matched += 1

    # 处理封面
    process_cover(slides[COVER_SLIDE], outline, prs, product_line)

    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    print(f"  处理页面: {matched}")


def detect_stage(md_path: str) -> str:
    """从文件名解析阶段"""
    name = Path(md_path).stem.upper()
    for key, val in {"CHARTER": "charter", "PDCP": "pdcp", "ADCP": "adcp",
                     "TRANSFER": "transfer", "转移": "transfer"}.items():
        if key in name:
            return val
    return ""


def main():
    parser = argparse.ArgumentParser(description="IPD 项目评审 PPT 生成（基于模板填充内容）")
    parser.add_argument("--outline", required=True, help="确认版大纲 Markdown 文件路径")
    parser.add_argument("--product-line", required=True, choices=["hm", "ml"],
                        help="产品线：hm（虹美）或 ml（美菱）")
    parser.add_argument("--output", default=".", help="输出目录")
    parser.add_argument("--stage", default=None,
                        help="阶段（可选，默认从文件名解析）")
    args = parser.parse_args()

    if not os.path.exists(args.outline):
        print(f"错误：大纲文件不存在: {args.outline}", file=sys.stderr)
        sys.exit(1)

    stage = args.stage or detect_stage(args.outline)
    if not stage:
        print("错误：无法从文件名解析阶段，请用 --stage 指定", file=sys.stderr)
        sys.exit(1)

    outline = parse_outline(args.outline)
    print(f"大纲解析完成: {len(outline['sections'])} 个章节, "
          f"{len(outline['notes_map'])} 条备注")

    template_path = resolve_template(args.product_line, stage)

    output_dir = Path(args.output)
    output_dir.mkdir(parents=True, exist_ok=True)

    outline_stem = Path(args.outline).stem
    output_path = str(output_dir / f"{outline_stem}-汇报材料.pptx")

    generate_ppt(outline, template_path, output_path, args.product_line)


if __name__ == "__main__":
    main()
