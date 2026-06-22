#!/usr/bin/env python3
"""
图表渲染引擎 v6 — 匹配模板比例

模板参数：
- 页面 960x540pt
- 章节标题 28pt，子标题 24pt
- 内容区域：x=53 y=124 w=854 h=416pt

字体比例：
- 图表标题：16-18pt（模板标题 28pt 的 60-65%）
- 正文/标签：13-14pt
- 小标签/序号：11-12pt
"""

import re
from dataclasses import dataclass
from enum import Enum

from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ══════════════════════════════════════════════════════
#  色板
# ══════════════════════════════════════════════════════

class C:
    BLUE = RGBColor(0x1E, 0x40, 0xAF)
    BLUE_L = RGBColor(0xDB, 0xEA, 0xFE)
    ORANGE = RGBColor(0xEA, 0x58, 0x0C)
    ORANGE_L = RGBColor(0xFF, 0xED, 0xD5)
    GREEN = RGBColor(0x16, 0xA3, 0x4A)
    GREEN_L = RGBColor(0xDC, 0xFC, 0xE7)
    PURPLE = RGBColor(0x7C, 0x3A, 0xED)
    PURPLE_L = RGBColor(0xED, 0xE9, 0xFE)
    CYAN = RGBColor(0x08, 0x91, 0xB2)
    CYAN_L = RGBColor(0xEC, 0xFE, 0xFF)
    PINK = RGBColor(0xDB, 0x27, 0x77)
    PINK_L = RGBColor(0xFC, 0xE7, 0xF3)
    AMBER = RGBColor(0xD9, 0x77, 0x06)
    AMBER_L = RGBColor(0xFE, 0xF3, 0xC7)

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    G50 = RGBColor(0xF9, 0xFA, 0xFB)
    G100 = RGBColor(0xF3, 0xF4, 0xF6)
    G200 = RGBColor(0xE5, 0xE7, 0xEB)
    G300 = RGBColor(0xD1, 0xD5, 0xDB)
    G500 = RGBColor(0x6B, 0x72, 0x80)
    G700 = RGBColor(0x37, 0x41, 0x51)
    G900 = RGBColor(0x11, 0x18, 0x27)

    PALETTE = [
        (BLUE, BLUE_L), (ORANGE, ORANGE_L), (GREEN, GREEN_L),
        (PURPLE, PURPLE_L), (CYAN, CYAN_L), (PINK, PINK_L),
        (AMBER, AMBER_L),
    ]


# ══════════════════════════════════════════════════════
#  图表类型
# ══════════════════════════════════════════════════════

class ChartType(str, Enum):
    PYRAMID = "pyramid"
    FLOW = "flow"
    COMPARISON = "comparison"
    HIERARCHY = "hierarchy"
    TIMELINE = "timeline"
    COMPOSITION = "composition"
    TABLE = "table"
    CARDS = "cards"


# ══════════════════════════════════════════════════════
#  内容分析
# ══════════════════════════════════════════════════════

def infer_chart_type(points, subtitle="", chart_hint=""):
    if chart_hint:
        h = chart_hint.lower().strip()
        if h in {ct.value for ct in ChartType}:
            return h
    clean = [p for p in points if not p.strip().startswith("<!--")]
    text = subtitle + " " + " ".join(clean)
    n = len(clean)
    if _tbl(clean): return "table"
    if _flow(text): return "flow"
    if _time(text): return "timeline"
    if _cmp(clean, subtitle): return "comparison"
    if _hier(clean, subtitle): return "hierarchy"
    if _comp(text): return "composition"
    if 3 <= n <= 6: return "pyramid"
    return "cards"

def _tbl(pts):
    return sum(1 for p in pts if "|" in p and not p.strip().startswith("<!--")) >= 2
def _flow(t):
    kws = ["→", "阶段", "步骤", "流程", "首先", "然后", "最后"]
    if sum(t.count(k) for k in kws) >= 2: return True
    return len(re.findall(r"(?:^|\s)\d+[.、）)]\s", t)) >= 2
def _time(t):
    if len(re.findall(r"20\d{2}[.\-/]\d{1,2}", t)) >= 2: return True
    node_kws = ["节点", "里程碑", "阶段一", "阶段二", "阶段三"]
    if sum(t.count(k) for k in node_kws) >= 2: return True
    return sum(t.count(k) for k in ["Q1", "Q2", "Q3", "Q4"]) >= 2
def _cmp(pts, sub):
    if any(k in sub for k in ["对比", "竞品", "方案", "vs"]): return True
    t = " ".join(pts)
    return sum(1 for k in ["我方", "竞品", "方案A", "方案B", "对比", "优于"] if k in t) >= 2
def _hier(pts, sub):
    if any(k in sub for k in ["架构", "框架", "体系", "组织", "结构"]): return True
    flow_kws = ["阶段", "步骤", "首先", "然后"]
    cols = [p for p in pts if ":" in p or "：" in p]
    if sum(1 for p in cols if any(k in p for k in flow_kws)) >= 2: return False
    return len(cols) >= 3
def _comp(t):
    return "%" in t and any(k in t for k in ["占比", "比例", "构成", "份额", "总计", "其中"])


# ══════════════════════════════════════════════════════
#  数据结构
# ══════════════════════════════════════════════════════

@dataclass
class Rect:
    x: int; y: int; w: int; h: int


# ══════════════════════════════════════════════════════
#  绘图基础
# ══════════════════════════════════════════════════════

def _s(slide, st, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(st, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def _t(slide, x, y, w, h, txt, sz=13, bold=False, color=None,
       align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    if color is None: color = C.G900
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.alignment = align; p.text = txt
    for r in p.runs:
        r.font.size = Pt(sz); r.font.name = font
        r.font.bold = bold; r.font.color.rgb = color
    return tb

def _trunc(t, n=20):
    return t[:n] + "…" if len(t) > n else t

def _parse(p):
    if "：" in p:
        k, v = p.split("：", 1); return k.strip(), v.strip()
    if ":" in p:
        k, v = p.split(":", 1); return k.strip(), v.strip()
    return p.strip(), ""


# ══════════════════════════════════════════════════════
#  渲染入口
# ══════════════════════════════════════════════════════

def render_chart(slide, chart_type, points, area, subtitle=""):
    fn = {
        "pyramid": _r_pyramid, "flow": _r_flow, "comparison": _r_compare,
        "hierarchy": _r_hierarchy, "timeline": _r_timeline,
        "composition": _r_composition, "table": _r_table, "cards": _r_cards,
    }.get(chart_type, _r_cards)
    fn(slide, points, area, subtitle)

def analyze_and_render(slide, points, area, subtitle="", hint=""):
    ct = infer_chart_type(points, subtitle, hint)
    render_chart(slide, ct, points, area, subtitle)
    return ct


# ══════════════════════════════════════════════════════
#  渲染器 v6 — 匹配模板比例
#  模板：960x540pt，内容区 854x416pt
#  字体：标题 16-18pt，正文 13-14pt，标签 11-12pt
# ══════════════════════════════════════════════════════

def _r_pyramid(slide, pts, area, sub=""):
    """
    金字塔：填满内容区，每层高 = area.h/n，字体 16pt
    """
    items = pts[:5]
    n = len(items)
    if not n: return

    layer_h = (area.h - Pt(10) * (n - 1)) // n
    max_w = area.w - Pt(20)

    for i, txt in enumerate(items):
        ratio = 1.0 - i * 0.08
        lw = int(max_w * ratio)
        lx = area.x + (area.w - lw) // 2
        ly = area.y + i * (layer_h + Pt(10))
        fg = C.PALETTE[i % len(C.PALETTE)][0]

        _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE, lx, ly, lw, layer_h, fill=fg)
        # 序号
        ns = Pt(30)
        _s(slide, MSO_SHAPE.OVAL, lx + Pt(14), ly + (layer_h - ns) // 2, ns, ns, fill=C.WHITE)
        _t(slide, lx + Pt(14), ly + (layer_h - ns) // 2, ns, ns,
           str(i+1), sz=14, bold=True, color=fg, align=PP_ALIGN.CENTER, font="Arial")
        # 文字（16pt）
        _t(slide, lx + Pt(50), ly + Pt(4), lw - Pt(60), layer_h - Pt(8),
           txt, sz=16, bold=True, color=C.WHITE)


def _r_flow(slide, pts, area, sub=""):
    """
    流程图：填满宽度，节点高度占内容区 60%，字体 15pt
    """
    items = pts[:5]
    n = len(items)
    if not n: return

    arrow_w = Pt(28)
    total_arrows = (n - 1) * arrow_w
    node_w = (area.w - total_arrows - Pt(20)) // n
    node_h = int(area.h * 0.55)
    start_x = area.x + Pt(10)
    start_y = area.y + (area.h - node_h) // 2

    for i, txt in enumerate(items):
        x = start_x + i * (node_w + arrow_w)
        fg = C.PALETTE[i % len(C.PALETTE)][0]

        _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, node_w, node_h, fill=fg)
        # 序号
        ns = Pt(30)
        _s(slide, MSO_SHAPE.OVAL, x + node_w//2 - ns//2, start_y - Pt(14), ns, ns, fill=C.WHITE, line=fg)
        _t(slide, x + node_w//2 - ns//2, start_y - Pt(14), ns, ns,
           str(i+1), sz=13, bold=True, color=fg, align=PP_ALIGN.CENTER, font="Arial")
        # 文字（15pt，不截断）
        _t(slide, x + Pt(10), start_y + Pt(12), node_w - Pt(20), node_h - Pt(16),
           txt, sz=15, bold=True, color=C.WHITE)

        if i < n - 1:
            ax = x + node_w + Pt(4)
            _t(slide, ax, start_y + node_h//2 - Pt(12), arrow_w - Pt(8), Pt(24),
               "→", sz=20, bold=True, color=C.G500)


def _r_compare(slide, pts, area, sub=""):
    """
    对比图：左右两列，填满内容区，字体 14pt
    """
    lg, rg = [], []
    for p in pts:
        if any(k in p for k in ["我方", "本项目", "方案A", "现有"]): lg.append(p)
        elif any(k in p for k in ["竞品", "方案B", "行业"]): rg.append(p)
        else: (lg if len(lg) <= len(rg) else rg).append(p)

    groups = [("我方 / 方案A", lg, C.BLUE), ("竞品 / 方案B", rg, C.ORANGE)]
    col_w = (area.w - Pt(20)) // 2
    bar_h = min(Pt(50), (area.h - Pt(80)) // 4)
    label_w = Pt(100)

    for ci, (gn, items, gc) in enumerate(groups):
        if not items: continue
        cx = area.x + ci * (col_w + Pt(20))

        _t(slide, cx, area.y, col_w, Pt(28),
           gn, sz=18, bold=True, color=gc)
        _s(slide, MSO_SHAPE.RECTANGLE, cx, area.y + Pt(30), col_w, Pt(3), fill=gc)

        y = area.y + Pt(40)
        for item in items[:4]:
            label, detail = _parse(item)
            _t(slide, cx, y, label_w, bar_h, label, sz=13, color=C.G700, align=PP_ALIGN.RIGHT)
            bw = col_w - label_w - Pt(10)
            _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + label_w + Pt(6), y + Pt(4), bw, bar_h - Pt(8), fill=gc)
            if detail:
                _t(slide, cx + label_w + Pt(12), y + Pt(4), bw - Pt(12), bar_h - Pt(8),
                   detail, sz=13, color=C.WHITE)
            y += bar_h + Pt(6)


def _r_hierarchy(slide, pts, area, sub=""):
    """
    层级图：撑满内容区宽度，每行高自适应，字体 14pt
    """
    items = pts[:7]
    n = len(items)
    if not n: return

    item_h = (area.h - Pt(6) * (n - 1)) // n
    item_h = min(item_h, Pt(56))

    for i, txt in enumerate(items):
        label, detail = _parse(txt)
        il = 0 if detail else 1
        indent = Pt(30) * il

        x = area.x + indent
        w = area.w - indent
        y = area.y + i * (item_h + Pt(6))
        fg, bg = C.PALETTE[i % len(C.PALETTE)]

        _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, item_h, fill=bg, line=C.G200)
        _s(slide, MSO_SHAPE.RECTANGLE, x, y + Pt(4), Pt(5), item_h - Pt(8), fill=fg)
        display = f"{label}：{detail}" if detail else label
        _t(slide, x + Pt(14), y + Pt(4), w - Pt(18), item_h - Pt(8),
           display, sz=14, bold=True, color=C.G900)


def _r_timeline(slide, pts, area, sub=""):
    """
    时间轴：圆点直径 28pt，标签卡片 130x50pt，字体 13pt
    """
    items = pts[:6]
    n = len(items)
    if not n: return

    ly = area.y + area.h // 2
    ll = area.x + Pt(20)
    lw = area.w - Pt(40)

    _s(slide, MSO_SHAPE.RECTANGLE, ll, ly - Pt(2), lw, Pt(4), fill=C.G300)

    gap = lw // max(n, 1)
    for i, txt in enumerate(items):
        x = ll + i * gap + gap // 2
        fg, bg = C.PALETTE[i % len(C.PALETTE)]

        ds = Pt(28)
        _s(slide, MSO_SHAPE.OVAL, x - ds//2, ly - ds//2, ds, ds, fill=fg)
        _t(slide, x - ds//2, ly - ds//2 + Pt(2), ds, ds - Pt(4),
           str(i+1), sz=12, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER, font="Arial")

        is_top = (i % 2 == 0)
        card_h = Pt(50)
        card_w = Pt(130)
        ct = ly - ds - Pt(6) - card_h if is_top else ly + ds + Pt(6)

        _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x - card_w//2, ct, card_w, card_h, fill=bg, line=C.G200)
        _s(slide, MSO_SHAPE.RECTANGLE, x - Pt(1), ct + card_h if is_top else ly + ds//2,
           Pt(2), ds + Pt(4) if is_top else Pt(6), fill=fg)
        _t(slide, x - card_w//2 + Pt(6), ct + Pt(6), card_w - Pt(12), card_h - Pt(12),
           txt, sz=13, bold=True, color=C.G900)


def _r_composition(slide, pts, area, sub=""):
    """
    占比图：撑满内容区宽度，条高自适应，字体 14pt
    """
    items = pts[:6]
    n = len(items)
    if not n: return

    bar_h = (area.h - Pt(8) * (n - 1)) // n
    bar_h = min(bar_h, Pt(52))

    # 标签占 20%，条形占 80%
    label_w = int(area.w * 0.20)
    pct_area_w = area.w - label_w - Pt(10)

    for i, txt in enumerate(items):
        y = area.y + i * (bar_h + Pt(8))
        fg, bg = C.PALETTE[i % len(C.PALETTE)]

        pct_m = re.search(r"(\d+(?:\.\d+)?)\s*%", txt)
        if pct_m:
            pct = float(pct_m.group(1)) / 100.0
            label = re.sub(r"\d+(?:\.\d+)?\s*%", "", txt).strip().strip("（）()：: ")
        else:
            pct, label = 1.0 / n, txt

        bw = max(Pt(20), int(pct_area_w * pct))

        _t(slide, area.x, y, label_w - Pt(8), bar_h,
           label, sz=14, color=C.G700, align=PP_ALIGN.RIGHT)
        _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
           area.x + label_w, y + Pt(3), bw, bar_h - Pt(6), fill=fg)
        if pct_m:
            _t(slide, area.x + label_w + bw + Pt(8), y, Pt(44), bar_h,
               f"{pct*100:.0f}%", sz=14, bold=True, color=fg, font="Arial")


def _r_table(slide, pts, area, sub=""):
    """美化表格：行高 28pt，字体 11-12pt"""
    rows = [p for p in pts if "|" in p]
    if not rows: return
    td = [[c.strip() for c in r.strip("|").split("|")] for r in rows]
    if not td: return

    nr, nc = len(td), max(len(r) for r in td)
    rh = min(Pt(28), (area.h - Pt(8)) // nr)
    cw = (area.w - Pt(8)) // nc

    for r, rd in enumerate(td):
        for c, ct in enumerate(rd):
            if c >= nc: break
            x = area.x + Pt(4) + c * cw
            y = area.y + Pt(4) + r * rh

            if r == 0: bg, tc, b, fs = C.BLUE, C.WHITE, True, 12
            else: bg = C.WHITE if r % 2 == 0 else C.G50; tc, b, fs = C.G900, False, 11

            s = _s(slide, MSO_SHAPE.RECTANGLE, x, y, cw - Pt(1), rh - Pt(1), fill=bg)
            s.line.color.rgb = C.G200; s.line.width = Pt(0.5)
            _t(slide, x + Pt(4), y + Pt(2), cw - Pt(8), rh - Pt(4),
               _trunc(ct, 18), sz=fs, bold=b, color=tc)


def _r_cards(slide, pts, area, sub=""):
    """
    卡片布局 — 填满内容区，卡片等分，字体 14pt
    """
    items = pts[:9]
    n = len(items)
    if not n: return

    if n <= 2:
        cols, rows = n, 1
    elif n <= 4:
        cols, rows = 2, 2
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 3, 3

    gap_x, gap_y = Pt(10), Pt(8)
    card_w = (area.w - gap_x * (cols + 1)) // cols
    card_h = (area.h - gap_y * (rows + 1)) // rows

    for i, txt in enumerate(items[:9]):
        r, c = i // cols, i % cols
        x = area.x + gap_x + c * (card_w + gap_x)
        y = area.y + gap_y + r * (card_h + gap_y)
        fg, bg = C.PALETTE[i % len(C.PALETTE)]

        _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, fill=C.WHITE, line=C.G200)
        _s(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, Pt(5), fill=fg)
        ns = Pt(28)
        _s(slide, MSO_SHAPE.OVAL, x + Pt(12), y + Pt(16), ns, ns, fill=fg)
        _t(slide, x + Pt(12), y + Pt(16), ns, ns,
           str(i+1), sz=13, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER, font="Arial")
        _t(slide, x + Pt(46), y + Pt(16), card_w - Pt(56), card_h - Pt(20),
           txt, sz=14, color=C.G900)
