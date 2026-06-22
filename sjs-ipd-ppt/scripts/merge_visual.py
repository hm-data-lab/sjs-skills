#!/usr/bin/env python3
"""
合并脚本：将视觉内容（PNG 图片）合并到基础 PPT 的对应页面

用法：
    python merge_visual.py <base_pptx> <visual_dir> [--output <output>]

base_pptx: 基础 PPT（gen_ppt.py 生成的）
visual_dir: 视觉内容目录，每张 PNG 文件名格式: slide_{N}.png (N=2,3,4...)
output: 输出路径，默认覆盖原文件
"""

import sys
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu


def find_visual_files(visual_dir):
    """
    扫描视觉内容目录，找到所有 slide_{N}.png 文件。
    返回 {slide_index: file_path} 字典。
    """
    files = {}
    visual_path = Path(visual_dir)
    for f in visual_path.glob("slide_*.png"):
        match = re.search(r"slide_(\d+)\.png", f.name)
        if match:
            slide_idx = int(match.group(1))
            files[slide_idx] = str(f)
    return files


def merge_visuals(base_pptx, visual_dir, output=None):
    """
    将视觉内容合并到基础 PPT。

    策略：
    - 找到 visual_dir 中的 slide_{N}.png
    - 在基础 PPT 的第 N 页（0-indexed: N-1）添加图片
    - 图片覆盖整个内容区域
    """
    if output is None:
        output = base_pptx

    prs = Presentation(base_pptx)
    visual_files = find_visual_files(visual_dir)

    if not visual_files:
        print(f"⚠ 未找到视觉文件: {visual_dir}/slide_*.png")
        return

    print(f"找到 {len(visual_files)} 个视觉文件")

    for slide_idx, file_path in sorted(visual_files.items()):
        # slide_N.png → slide index N-1 (0-indexed)
        prs_idx = slide_idx - 1
        if prs_idx < 0 or prs_idx >= len(prs.slides):
            print(f"  ⚠ slide_{slide_idx}.png: 超出范围 (共 {len(prs.slides)} 页)")
            continue

        slide = prs.slides[prs_idx]

        # 找到内容区域（跳过标题 shape[0] 和子标题 shape[1]）
        shapes = list(slide.shapes)
        content_top = Pt(124)  # 默认
        if len(shapes) > 1 and shapes[1].has_text_frame:
            content_top = shapes[1].top + shapes[1].height + Pt(5)

        # 内容区域参数
        content_left = Pt(53)
        content_width = prs.slide_width - Pt(106)  # 左右各留 53pt
        content_height = prs.slide_height - content_top - Pt(20)  # 底部留 20pt

        # 清空现有内容 shape（保留标题 shape[0] 和子标题 shape[1]）
        for shape in list(slide.shapes)[2:]:
            sp = shape._element
            sp.getparent().remove(sp)

        # 添加图片
        slide.shapes.add_picture(
            file_path,
            content_left, content_top,
            content_width, content_height,
        )
        print(f"  ✅ slide_{slide_idx}: {Path(file_path).name} → 内容区域")

    prs.save(output)
    print(f"\n✅ 合并完成: {output}")


def main():
    if len(sys.argv) < 3:
        print("用法: python merge_visual.py <base_pptx> <visual_dir> [--output <output>]")
        print("  base_pptx: 基础 PPT")
        print("  visual_dir: 视觉内容目录 (含 slide_2.png, slide_3.png ...)")
        sys.exit(1)

    base_pptx = sys.argv[1]
    visual_dir = sys.argv[2]
    output = None

    if "--output" in sys.argv:
        idx = sys.argv.index("--output")
        if idx + 1 < len(sys.argv):
            output = sys.argv[idx + 1]

    if not os.path.exists(base_pptx):
        print(f"错误: 基础 PPT 不存在: {base_pptx}")
        sys.exit(1)

    if not os.path.isdir(visual_dir):
        print(f"错误: 视觉内容目录不存在: {visual_dir}")
        sys.exit(1)

    merge_visuals(base_pptx, visual_dir, output)


if __name__ == "__main__":
    main()
