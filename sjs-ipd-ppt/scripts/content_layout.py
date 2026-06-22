#!/usr/bin/env python3
"""
内容布局增强（集成 chart_engine 自动图表推断）

原则：不改背景，只在模板上添加清晰的内容布局
- 自动推断每页内容适合的图表类型（金字塔、流程图、对比图等）
- 也保留手动标注 <!-- chart: xxx --> 的覆盖机制
- 关键数字突出，页面填满不留空白

支持模式：
  auto    — 自动推断图表类型（默认）
  manual  — 只用 <!-- chart: xxx --> 标注，未标注的用默认卡片
  none    — 跳过图表增强
"""

import sys
import re
from pathlib import Path

from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from chart_engine import (
    infer_chart_type, render_chart, Rect,
    C, ChartType,
    _s, _t,
)
from outline_parser import parse_outline_simple


def _to_chart_hint(slides):
    """将 outline_parser 格式转换为 content_layout 需要的格式（带 chart_hint）"""
    result = {}
    for heading, sub_list in slides.items():
        result[heading] = []
        for sub in sub_list:
            result[heading].append({
                "subtitle": sub.get("subtitle"),
                "points": sub.get("points", []),
                "chart_hint": sub.get("chart", ""),
            })
    return result


# ── 页面增强函数 ──

def _get_slide_area(slide, prs):
    """
    获取 slide 内容区域 — 匹配模板布局

    模板实际参数（pt）：
    - 页面：960 x 540
    - 子标题 shape[1]：x=53, y=79, w=194, h=45
    - 内容区域：x=53, y=124(=79+45), 到页面底部

    返回值坐标单位：EMU (1pt = 12700 EMU)
    """
    # 固定匹配模板内容区
    margin_x = Pt(53)
    margin_bottom = Pt(20)

    # 找子标题底部（shape[1]）
    shapes = list(slide.shapes)
    content_top = Pt(124)  # 默认：子标题底部
    if len(shapes) > 1 and shapes[1].has_text_frame:
        content_top = shapes[1].top + shapes[1].height + Pt(5)

    content_w = prs.slide_width - margin_x * 2
    content_h = prs.slide_height - content_top - margin_bottom

    return Rect(
        x=int(margin_x),
        y=int(content_top),
        w=int(content_w),
        h=int(content_h),
    )


def _find_content_points_for_slide(outline, subtitle_t):
    """根据 slide 的子标题文字找到对应的大纲要点"""
    subtitle_norm = re.sub(r"\s+", "", subtitle_t.strip())

    for heading, subs in outline.items():
        for sub in subs:
            sub_title = sub.get("subtitle", "")
            if sub_title and re.sub(r"\s+", "", sub_title) in subtitle_norm:
                return sub.get("points", []), sub.get("chart_hint", "")
            if sub_title and subtitle_norm in re.sub(r"\s+", "", sub_title):
                return sub.get("points", []), sub.get("chart_hint", "")

    return [], ""


def enhance_content_pages_auto(prs, outline, stage="charter"):
    """
    内容页：混合方案自动推断。

    复杂图表（pyramid/flow/comparison/timeline）→ SVG 模板生成
    简单布局（cards/hierarchy/composition/table）→ python-pptx shapes
    """
    from utils import get_stage_config
    cfg = get_stage_config(stage)
    content_start, content_end = cfg["content_range"]

    # 尝试导入 SVG 图表生成器
    try:
        from svg_chart_gen import render_pyramid, render_flow, render_comparison, render_timeline
        has_svg = True
    except ImportError:
        has_svg = False

    for i in range(content_start, content_end):
        if i >= len(prs.slides):
            break

        slide = prs.slides[i]
        shapes = list(slide.shapes)
        if len(shapes) < 2:
            continue

        subtitle = ""
        if shapes[1].has_text_frame:
            subtitle = shapes[1].text_frame.text.strip().split("\n")[0]

        points, chart_hint = _find_content_points_for_slide(outline, subtitle)
        if not points:
            continue

        # 保存原文到备注
        if len(shapes) > 2 and shapes[2].has_text_frame:
            original_t = shapes[2].text_frame.text.strip()
            if original_t:
                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                slide.notes_slide.notes_text_frame.text = (notes + "\n\n【原始内容】\n" + original_t).strip()

        # 清空正文
        if len(shapes) > 2 and shapes[2].has_text_frame:
            tf = shapes[2].text_frame
            while len(tf.paragraphs) > 1:
                p = tf.paragraphs[-1]
                p._element.getparent().remove(p._element)
            tf.paragraphs[0].text = ""

        area = _get_slide_area(slide, prs)
        if not area:
            continue

        chart_type = infer_chart_type(points, subtitle, chart_hint)

        # 混合方案：复杂图表用 SVG，简单布局用 shapes
        if has_svg and chart_type in ("pyramid", "flow", "comparison", "timeline"):
            if chart_type == "pyramid":
                render_pyramid(slide, points, area, prs, title=subtitle)
            elif chart_type == "flow":
                render_flow(slide, points, area, prs, title=subtitle)
            elif chart_type == "comparison":
                # 分左右组
                lg = [p for p in points if any(k in p for k in ["我方", "本项目", "方案A"])]
                rg = [p for p in points if any(k in p for k in ["竞品", "方案B", "行业"])]
                if not lg: lg = points[:len(points)//2]
                if not rg: rg = points[len(points)//2:]
                render_comparison(slide, lg, rg, area, prs,
                                  left_title="我方", right_title="竞品", title=subtitle)
            elif chart_type == "timeline":
                render_timeline(slide, points, area, prs, title=subtitle)
            print(f"  Slide {i + 1}: {subtitle[:20]} → {chart_type} (SVG)")
        else:
            render_chart(slide, chart_type, points, area, subtitle)
            print(f"  Slide {i + 1}: {subtitle[:20]} → {chart_type}")


def enhance_value_page(prs, outline, stage="charter"):
    """
    项目价值页：四维度卡片 / 雷达图。

    优先用 <!-- chart: xxx --> 标注，否则用默认 2×2 四象限。
    """
    from utils import get_stage_config
    cfg = get_stage_config(stage)
    value_page = cfg["value_page"]
    if value_page is None or value_page >= len(prs.slides):
        return

    slide = prs.slides[value_page]
    value_subs = outline.get("三、项目价值", [])
    if not value_subs:
        return

    # 清空现有内容（保留标题）
    shapes = list(slide.shapes)
    for shape in shapes[1:]:
        sp = shape._element
        sp.getparent().remove(sp)

    # 检查是否有手动标注
    manual_hint = ""
    for sub in value_subs:
        if sub.get("chart_hint"):
            manual_hint = sub["chart_hint"]
            break

    if manual_hint and manual_hint != "radar":
        # 用指定的图表类型
        all_points = []
        for sub in value_subs:
            title = sub.get("subtitle", "")
            pts = sub.get("points", [])
            if title:
                all_points.append(f"{title}: {'; '.join(pts[:2])}")
            else:
                all_points.extend(pts)

        area = Rect(
            x=Emu(50000), y=Emu(110000),
            w=prs.slide_width - Emu(100000),
            h=prs.slide_height - Emu(160000),
        )
        render_chart(slide, manual_hint, all_points, area)
    else:
        # 默认 2×2 四象限卡片
        card_w = Pt(420)
        card_h = Pt(170)
        gap_x = Pt(16)
        gap_y = Pt(14)
        start_x = Pt(53)
        start_y = Pt(110)
        card_colors = [C.BLUE, C.ORANGE, C.CYAN, C.GREEN]

        for idx, sub in enumerate(value_subs[:4]):
            row = idx // 2
            col = idx % 2
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)
            title = sub.get("subtitle", "")
            points = sub.get("points", [])
            if title:
                _render_value_card(slide, x, y, card_w, card_h, title, points, card_colors[idx])

    print("  Slide 9: 项目价值 → 四维度卡片")


def _render_value_card(slide, left, top, width, height, title, points, color):
    """渲染单个价值维度卡片"""
    # 背景
    _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
              left, top, width, height, fill=C.WHITE)

    # 顶部色条
    _s(slide, MSO_SHAPE.RECTANGLE,
              left, top, width, Pt(4), fill=color)

    # 标题
    _t(slide, left + Pt(12), top + Pt(10), width - Pt(24), Pt(22),
                  title, sz=14, bold=True, color=color, align=PP_ALIGN.LEFT)

    # 要点
    py = top + Pt(38)
    for pt in points[:4]:
        _t(slide, left + Pt(12), py, width - Pt(24), Pt(18),
                      f"• {pt}", sz=11, color=C.G700, align=PP_ALIGN.LEFT)
        py += Pt(20)


def enhance_kpi_page(prs, outline, stage="charter"):
    """产出收益页：KPI 大数字卡片"""
    from utils import get_stage_config
    cfg = get_stage_config(stage)
    output_page = cfg["output_page"]
    if output_page is None or output_page >= len(prs.slides):
        return

    slide = prs.slides[output_page]
    kpi_points = []
    for sub in outline.get("七、产出和收益", []):
        for pt in sub.get("points", []):
            if "⭐" in pt:
                kpi_points.append(pt)

    if not kpi_points:
        return

    # 清空现有内容（保留标题）
    shapes = list(slide.shapes)
    for shape in shapes[1:]:
        sp = shape._element
        sp.getparent().remove(sp)

    # 解析 KPI 数据
    kpis = []
    for pt in kpi_points[:4]:
        match = re.search(r"\*\*(.+?)\*\*[：:](.+)", pt)
        if match:
            kpis.append((match.group(1), match.group(2)[:20]))
        else:
            kpis.append((pt[:15], ""))

    # 渲染 KPI 卡片
    card_w = Pt(210)
    card_h = Pt(150)
    gap = Pt(16)
    start_x = Pt(53)
    start_y = Pt(130)
    kpi_colors = [C.ORANGE, C.BLUE, C.CYAN, C.GREEN]

    for idx, (number, label) in enumerate(kpis):
        x = start_x + idx * (card_w + gap)
        _render_kpi_card(slide, x, start_y, card_w, card_h, number, label, kpi_colors[idx])

    print("  Slide 15: 产出收益 → KPI 卡片")


def _render_kpi_card(slide, left, top, width, height, number, label, color):
    """渲染单个 KPI 卡片"""
    # 背景
    _s(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
              left, top, width, height, fill=C.WHITE)

    # 顶部色条
    _s(slide, MSO_SHAPE.RECTANGLE,
              left, top, width, Pt(5), fill=color)

    # 大数字
    _t(slide, left, top + Pt(20), width, Pt(45),
                  number, sz=30, bold=True, color=color,
                  font="Arial")

    # 标签
    _t(slide, left + Pt(8), top + Pt(70), width - Pt(16), Pt(40),
                  label, sz=11, color=C.G700)


# ── 主逻辑 ──

def enhance_ppt(pptx_path, outline_path=None, chart_mode="auto", stage="charter"):
    """
    增强 PPT 的内容布局。

    Args:
        pptx_path: PPT 文件路径
        outline_path: 大纲 Markdown 路径（可选）
        chart_mode: 图表模式 — auto / manual / none
        stage: 评审阶段 — charter / pdcp / adcp / transfer
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)

    outline = {}
    if outline_path and Path(outline_path).exists():
        outline = _to_chart_hint(parse_outline_simple(outline_path))

    if chart_mode == "none":
        print("图表增强已跳过（mode=none）")
        return prs

    from utils import get_stage_config
    cfg = get_stage_config(stage)
    content_start, content_end = cfg["content_range"]

    print(f"增强内容页（{content_start}-{content_end - 1}）[mode={chart_mode}]...")
    enhance_content_pages_auto(prs, outline, stage)

    print("增强项目价值页...")
    enhance_value_page(prs, outline, stage)

    print("增强产出收益页...")
    enhance_kpi_page(prs, outline, stage)

    prs.save(pptx_path)
    print(f"✅ 内容布局完成: {pptx_path}")
    return prs


def main():
    if len(sys.argv) < 2:
        print("用法: python content_layout.py <pptx> [outline.md] [--mode auto|manual|none]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    outline_path = None
    chart_mode = "auto"

    for i, arg in enumerate(sys.argv[2:], 2):
        if arg == "--mode" and i + 1 < len(sys.argv):
            chart_mode = sys.argv[i + 1]
        elif not arg.startswith("--"):
            outline_path = arg

    enhance_ppt(pptx_path, outline_path, chart_mode)


if __name__ == "__main__":
    main()
