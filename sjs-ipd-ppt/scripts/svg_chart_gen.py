#!/usr/bin/env python3
"""
SVG 图表生成器 — 基于 ppt-master 模板

将数据填充到 ppt-master 的 SVG 模板中，生成高质量图表图片。
用于混合方案：复杂图表用 SVG，简单布局用 python-pptx shapes。
"""

import os
import re
import tempfile
from pathlib import Path

import cairosvg
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt, Emu

from utils import get_template_dir, fix_chinese_font

TEMPLATE_DIR = get_template_dir()


# ══════════════════════════════════════════════════════
#  SVG 模板填充
# ══════════════════════════════════════════════════════

def _load_svg(name):
    """加载 SVG 模板"""
    path = TEMPLATE_DIR / f"{name}.svg"
    if not path.exists():
        raise FileNotFoundError(f"SVG template not found: {path}")
    return path.read_text(encoding="utf-8")


def _replace_text(svg, replacements):
    """
    替换 SVG 中的文本。
    replacements: list of (old_text, new_text)
    """
    for old, new in replacements:
        # XML 实体转义
        new_escaped = new.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        svg = svg.replace(old, new_escaped)
    return svg




def _svg_to_png(svg_content, width=1280, height=720):
    """SVG → PNG bytes"""
    png_data = cairosvg.svg2png(
        bytestring=svg_content.encode("utf-8"),
        output_width=width,
        output_height=height,
    )
    return png_data


def _png_to_slide(slide, png_bytes, area, prs):
    """将 PNG 嵌入 slide 的指定区域"""
    # 保存为临时文件
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        f.write(png_bytes)
        tmp_path = f.name

    try:
        # 添加图片
        slide.shapes.add_picture(tmp_path, area.x, area.y, area.w, area.h)
    finally:
        os.unlink(tmp_path)


# ══════════════════════════════════════════════════════
#  金字塔图
# ══════════════════════════════════════════════════════

def gen_pyramid(items, title="", subtitle=""):
    """
    生成金字塔图 SVG。

    items: list of str, 从底层到顶层（或从顶到底，自动反转）
    title: 图表标题
    subtitle: 副标题
    """
    svg = _load_svg("pyramid_chart")

    # 确保最多 5 层
    items = items[:5]
    # 金字塔从底层(L1)到顶层(L5)，但用户通常从顶层开始列
    # 所以反转：用户的第1项 = 金字塔顶层(L5)
    items = list(reversed(items))

    # 填充标题
    if not title:
        title = "Strategic Pyramid"
    if not subtitle:
        subtitle = "FIVE-TIER HIERARCHY"

    svg = svg.replace("Strategic Capability Pyramid", title)
    svg = svg.replace("FIVE-TIER EVOLUTION MODEL · FROM FOUNDATION TO VISION", subtitle)

    # 填充层级文字（从顶层 L5 到底层 L1）
    layer_labels = ["VISION", "INNOVATION", "BRAND POWER", "R&amp;D SYSTEM", "OPERATIONS"]
    for i, item in enumerate(items):
        if i < len(layer_labels):
            label = item[:18].upper()
            svg = svg.replace(layer_labels[i], label)

    # 填充右侧卡片文字
    card_titles = ["Vision &amp; Purpose", "Innovation Engine", "Brand Power", "R&amp;D System", "Operations Foundation"]
    card_descs = [
        "Establishes industry leadership and guides overall strategic layout.",
        "Drives technology innovation and builds core competitive advantages.",
        "Enhances brand influence and market recognition.",
        "Builds systematic research and development capabilities.",
        "Ensures efficient operational execution and delivery.",
    ]

    for i, item in enumerate(items):
        if i < len(card_titles):
            # 卡片标题用 items 的内容
            new_title = item[:30]
            svg = svg.replace(card_titles[i], new_title)
            # 卡片描述用 items 的内容
            if i < len(items):
                new_desc = item[:60] if len(item) > 30 else ""
                svg = svg.replace(card_descs[i], new_desc)

    return svg


def render_pyramid(slide, items, area, prs, title="", subtitle=""):
    """渲染金字塔图到 slide"""
    svg = gen_pyramid(items, title, subtitle)
    svg = fix_chinese_font(svg)
    png = _svg_to_png(svg)
    _png_to_slide(slide, png, area, prs)


# ══════════════════════════════════════════════════════
#  流程图
# ══════════════════════════════════════════════════════

def gen_flow(items, title="", subtitle=""):
    """生成流程图 SVG"""
    svg = _load_svg("process_flow")

    items = items[:6]

    if not title:
        title = "Process Flow"
    if not subtitle:
        subtitle = "SEQUENTIAL STEPS"

    svg = svg.replace("Process Flow", title)
    svg = svg.replace("SEQUENTIAL STEPS", subtitle)

    # 流程图模板有 step 标签，需要找到并替换
    # 查找所有 text 元素中的 step 标签
    step_patterns = re.findall(r'<text[^>]*>([^<]*(?:Step|STEP|Phase|PHASE)[^<]*)</text>', svg)
    for i, pattern in enumerate(step_patterns):
        if i < len(items):
            new_text = items[i][:20]
            svg = svg.replace(pattern, new_text, 1)

    return svg


def render_flow(slide, items, area, prs, title="", subtitle=""):
    """渲染流程图到 slide"""
    svg = gen_flow(items, title, subtitle)
    svg = fix_chinese_font(svg)
    png = _svg_to_png(svg)
    _png_to_slide(slide, png, area, prs)


# ══════════════════════════════════════════════════════
#  对比图
# ══════════════════════════════════════════════════════

def gen_comparison(left_items, right_items, left_title="方案A", right_title="方案B", title="对比分析"):
    """生成对比图 SVG"""
    svg = _load_svg("comparison_columns")

    svg = svg.replace("Comparison", title)

    # 填充左右列标题
    col_titles = re.findall(r'<text[^>]*class="col-title"[^>]*>([^<]*)</text>', svg)
    if not col_titles:
        col_titles = re.findall(r'<text[^>]*>([^<]*(?:Column|Team|Option)[^<]*)</text>', svg)

    if len(col_titles) >= 2:
        svg = svg.replace(col_titles[0], left_title)
        svg = svg.replace(col_titles[1], right_title)

    return svg


def render_comparison(slide, left_items, right_items, area, prs,
                      left_title="方案A", right_title="方案B", title="对比分析"):
    """渲染对比图到 slide"""
    svg = gen_comparison(left_items, right_items, left_title, right_title, title)
    svg = fix_chinese_font(svg)
    png = _svg_to_png(svg)
    _png_to_slide(slide, png, area, prs)


# ══════════════════════════════════════════════════════
#  时间轴
# ══════════════════════════════════════════════════════

def gen_timeline(items, title="", subtitle=""):
    """生成时间轴 SVG"""
    svg = _load_svg("timeline")

    items = items[:8]

    if not title:
        title = "Project Timeline"
    if not subtitle:
        subtitle = "KEY MILESTONES"

    svg = svg.replace("Timeline", title)
    svg = svg.replace("KEY MILESTONES", subtitle)

    return svg


def render_timeline(slide, items, area, prs, title="", subtitle=""):
    """渲染时间轴到 slide"""
    svg = gen_timeline(items, title, subtitle)
    svg = fix_chinese_font(svg)
    png = _svg_to_png(svg)
    _png_to_slide(slide, png, area, prs)


# ══════════════════════════════════════════════════════
#  KPI 卡片
# ══════════════════════════════════════════════════════

def gen_kpi_cards(kpis, title=""):
    """
    生成 KPI 卡片 SVG。
    kpis: list of (number, label)
    """
    svg = _load_svg("kpi_cards")

    if not title:
        title = "Key Performance Indicators"

    svg = svg.replace("Key Performance Indicators", title)

    return svg


def render_kpi_cards(slide, kpis, area, prs, title=""):
    """渲染 KPI 卡片到 slide"""
    svg = gen_kpi_cards(kpis, title)
    svg = fix_chinese_font(svg)
    png = _svg_to_png(svg)
    _png_to_slide(slide, png, area, prs)
